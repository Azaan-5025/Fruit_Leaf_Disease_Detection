import os
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
from collections import Counter
from PIL import Image
import cv2
import json

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ------------------------------------------------------------------
# Configuration
# ------------------------------------------------------------------
NAME = "Syed Azaan Hussain"
REG_NO = "24BCE5025"
PPTX_PATH = r"c:\Users\Syed Azaan Hussain\DS_Project\DS project Review 2.pptx"
DATA_DIR = r"c:\Users\Syed Azaan Hussain\DS_Project\data\training_data"
ASSETS_DIR = r"c:\Users\Syed Azaan Hussain\DS_Project\review_assets"
os.makedirs(ASSETS_DIR, exist_ok=True)

# Colors
BG_SLIDE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)

# ------------------------------------------------------------------
# Helper Functions
# ------------------------------------------------------------------
def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.33), Inches(2))
    tf = txBox.text_frame
    tf.text = "Deep Learning Based Hybrid System for\nFruit & Leaf Disease Detection"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.33), Inches(2))
    tf = txBox.text_frame
    tf.text = f"Presented by: {NAME}\nRegister No: {REG_NO}"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list=None, code_snippet=None, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT_BLUE; shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)

    if content_list:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for item in content_list:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.space_after = Pt(10)

    if code_snippet:
        txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(5))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = code_snippet
        p.font.size = Pt(10)
        p.font.name = 'Courier New'
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(1.4), Inches(6.2), Inches(5.2))
        shape.fill.background(); shape.line.color.rgb = ACCENT_BLUE
        # Re-order
        # No-op for z_order as default is correct

    if image_path:
        slide.shapes.add_picture(image_path, Inches(7), Inches(1.5), width=Inches(5.5))

# ------------------------------------------------------------------
# Main Logic
# ------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

add_title_slide(prs)

# Abstract
add_content_slide(prs, "Abstract", [
    "Development of a high-fidelity disease identification system using unified neural architectures.",
    "Specifically targets 30+ variants of leaf and fruit diseases with localized feature extraction.",
    "Utilizes the Sakaguchi preprocessing pipeline to enhance spatial importance in low-resolution tensors.",
    "Integrated formalin detection module for food safety assurance.",
    "Proposed hybrid model combines MobileNetV2 depth with auxiliary spatial feature maps."
])

# Introduction
add_content_slide(prs, "Introduction", [
    "Agrarian economics face significant threats from rapid disease spread and quality degradation.",
    "Traditional manual inspection is slow, subjective, and prone to human error.",
    "Modern solutions require edge-compatible Deep Learning models for real-time field diagnostics.",
    "The Sakaguchi method addresses noise sensitivity and varied illumination in field photography.",
    "Objective: Accuracy > 85% with inference latency < 3s."
])

# Dataset Detail
classes = sorted([d for d in os.listdir(DATA_DIR) if os.path.isdir(os.path.join(DATA_DIR, d))])
add_content_slide(prs, "Dataset Description", [
    f"Total Categories: {len(classes)} classes including healthy & diseased states.",
    "Domains: Apple, Corn, Grape, Tomato, Banana (Fresh/Rotten/Formalin).",
    "Source: Curated from PlantVillage and localized field datasets.",
    "Balance: Augmentation used to handle skewed class representations.",
    "Resolution: Multi-scale processing (224 to 8px)."
])

# EDA Code & Output
# Generate EDA Chart if not exists
plt.figure(figsize=(10, 6))
counts = [len(os.listdir(os.path.join(DATA_DIR, c))) for c in classes[:10]]
sns.barplot(x=classes[:10], y=counts, palette='viridis')
plt.xticks(rotation=45, ha='right', fontsize=8)
plt.title("Sample Class Distribution (Top 10)")
eda_path = os.path.join(ASSETS_DIR, "eda_review.png")
plt.savefig(eda_path, dpi=150, bbox_inches='tight')
plt.close()

eda_code = """
import os
import seaborn as sns
import matplotlib.pyplot as plt

def plot_distribution(data_path):
    classes = os.listdir(data_path)
    counts = [len(os.listdir(os.path.join(data_path, c))) 
              for c in classes]
    
    plt.figure(figsize=(12, 6))
    sns.barplot(x=classes, y=counts)
    plt.title("Distribution of Samples per Class")
    plt.show()
"""
add_content_slide(prs, "EDA: Code and Visual Analysis", [
    "Statistical audits ensure data consistency before training.",
    "Distribution analysis identifies minority classes requiring oversampling.",
    "Visual verification of image headers and pixel depth."
], code_snippet=eda_code, image_path=eda_path)

# Data Pre-Processing
pre_code = """
def sakaguchi_pipeline(img):
    # 1. Smoothing
    blur = cv2.GaussianBlur(img, (5,5), 0)
    # 2. Resizing
    res = cv2.resize(blur, (224,224))
    # 3. Tensor Extraction
    t8 = cv2.resize(res, (8,8))
    t12 = cv2.resize(res, (12,12))
    t16 = cv2.resize(res, (16,16))
    return t8, t12, t16
"""

# Generate Preprocessing Vizes
sample_class = classes[0]
sample_file = os.listdir(os.path.join(DATA_DIR, sample_class))[0]
img = cv2.imread(os.path.join(DATA_DIR, sample_class, sample_file))
img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
t8 = cv2.resize(img, (8,8))
t12 = cv2.resize(img, (12,12))
t16 = cv2.resize(img, (16,16))

fig, axes = plt.subplots(1, 3, figsize=(12, 4))
axes[0].imshow(t8); axes[0].set_title("8x8 Tensor")
axes[1].imshow(t12); axes[1].set_title("12x12 Tensor")
axes[2].imshow(t16); axes[2].set_title("16x16 Tensor")
pre_path = os.path.join(ASSETS_DIR, "pre_review.png")
plt.savefig(pre_path, dpi=150, bbox_inches='tight')
plt.close()

add_content_slide(prs, "Preprocessing: Sakaguchi Pipeline", [
    "Noise reduction via localized Gaussian Kernels.",
    "Spatial Compression: Downsampling to discrete tensor representations.",
    "Goal: Capture macroscopic structural patterns that are noise-invariant.",
    "The 8x8 configuration highlights global color distributions.",
    "The 16x16 configuration retains significant textural intensity."
], code_snippet=pre_code, image_path=pre_path)

# Methodology
add_content_slide(prs, "Methodology & Architecture", [
    "Model: Multi-Input Hybrid Convolutional Neural Network.",
    "Base Extractor: MobileNetV2 (Pre-trained on ImageNet).",
    "Auxiliary Channels: Three parallel Dense layers for Sakaguchi Tensors.",
    "Feature Fusion: Late-stage concatenation of CNN weights + Tensor features.",
    "Optimizer: Adam with exponential decay learning rate scheduling.",
    "Loss Function: Categorical Cross-Entropy for multi-class discriminators."
])

# Methodology Visual (Placeholder or explanation)
add_content_slide(prs, "Expected Output & Target Metrics", [
    "Primary Metric: Validation Accuracy targeting > 85%.",
    "Current Achievement: 85.05% Accuracy on Multi-Input CNN.",
    "Reliability: Integrated majority-voting consensus across 4 model types.",
    "Output: Real-time status (Healthy/Diseased) + Specific Disease Label.",
    "Confidence: Visual probability mapping for decision transparency."
])

prs.save(PPTX_PATH)
print(f"Presentation saved to: {PPTX_PATH}")
