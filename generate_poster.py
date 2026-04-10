"""
Generate a VIT-style academic poster for DS Project.
Format: 3-column layout with header, matching the provided template.
"""

import os
import json
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# CONFIGURATION
# ============================================================
STUDENT_NAME = "Syed Azaan Hussain"
REGISTER_NO  = "24BCE5025"
GUIDE_NAME   = "Pattabhiraman V"
SCHOOL_NAME  = "School of Computer Science and Engineering"
GITHUB_LINK  = "https://github.com/Azaan-5025/Fruit_Leaf_Disease_Detection"
OUTPUT_FILE  = "DS_Project_Poster.pptx"

# Colors matching VIT poster template
VIT_BLUE    = RGBColor(0x1A, 0x6F, 0xB5)
VIT_DARK    = RGBColor(0x0D, 0x47, 0xA1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
HEADING_BLUE = RGBColor(0x0D, 0x6E, 0xB8)
BORDER_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

# Poster dimensions (A1-ish landscape on a single slide)
POSTER_W = Inches(20)
POSTER_H = Inches(15)

ASSETS_DIR = "ppt_assets"
os.makedirs(ASSETS_DIR, exist_ok=True)


# ============================================================
# Helpers
# ============================================================
def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=11, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
             font_name="Arial", space_after=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.name = font_name
    p.space_after = space_after
    return p

def add_section_heading(slide, left, top, width, text):
    tf = add_textbox(slide, left, top, width, Inches(0.4), text,
                     font_size=16, bold=True, color=HEADING_BLUE, font_name="Arial")
    return tf

def add_bordered_box(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER_GRAY
    shape.line.width = Pt(1)
    return shape


# ============================================================
# Chart Generators
# ============================================================
def generate_model_accuracy_donut():
    """Donut charts for CNN and SVM accuracy."""
    fig, axes = plt.subplots(1, 2, figsize=(5, 2.5))
    fig.patch.set_facecolor('white')

    # CNN donut
    cnn_acc = 85.1
    axes[0].pie([cnn_acc, 100-cnn_acc], colors=['#0D6EB8', '#E0E0E0'],
                startangle=90, wedgeprops=dict(width=0.35))
    axes[0].text(0, 0, f'{cnn_acc}%', ha='center', va='center',
                 fontsize=14, fontweight='bold', color='#0D6EB8')
    axes[0].set_title('CNN', fontsize=11, fontweight='bold', color='#333')

    # SVM donut
    svm_acc = 82.3
    axes[1].pie([svm_acc, 100-svm_acc], colors=['#FF9800', '#E0E0E0'],
                startangle=90, wedgeprops=dict(width=0.35))
    axes[1].text(0, 0, f'{svm_acc}%', ha='center', va='center',
                 fontsize=14, fontweight='bold', color='#FF9800')
    axes[1].set_title('SVM', fontsize=11, fontweight='bold', color='#333')

    plt.tight_layout()
    path = os.path.join(ASSETS_DIR, "poster_donuts.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def generate_model_comparison_bars():
    """Bar chart comparing all 4 models."""
    results_path = "models/comparison_results.json"
    if os.path.exists(results_path):
        with open(results_path, "r") as f:
            data = json.load(f)
        names = [d['model_name'] for d in data]
        train_acc = [d['train_accuracy']*100 for d in data]
        val_acc = [d['val_accuracy']*100 for d in data]
    else:
        names = ["CNN", "SVM", "ANN", "ResNet50"]
        train_acc = [90.1, 84.2, 42.9, 29.1]
        val_acc = [85.1, 82.3, 39.7, 25.1]

    x = np.arange(len(names))
    width = 0.35

    fig, ax = plt.subplots(figsize=(5.5, 2.5))
    fig.patch.set_facecolor('white')
    ax.set_facecolor('white')

    bars1 = ax.bar(x - width/2, train_acc, width, label='Train', color='#0D6EB8', edgecolor='none')
    bars2 = ax.bar(x + width/2, val_acc, width, label='Validation', color='#4CAF50', edgecolor='none')

    for bar in bars1:
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                f'{bar.get_height():.1f}%', ha='center', va='bottom', fontsize=7, color='#333')
    for bar in bars2:
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                f'{bar.get_height():.1f}%', ha='center', va='bottom', fontsize=7, color='#333')

    ax.set_ylabel('Accuracy (%)', fontsize=9, color='#333')
    ax.set_xticks(x)
    ax.set_xticklabels(names, fontsize=8, color='#333')
    ax.legend(fontsize=8, frameon=False)
    ax.set_ylim(0, 110)
    for spine in ax.spines.values():
        spine.set_color('#CCC')
    ax.tick_params(colors='#333')
    plt.tight_layout()
    path = os.path.join(ASSETS_DIR, "poster_comparison.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def generate_architecture_diagram():
    """Architecture flow diagram for the poster."""
    fig, ax = plt.subplots(figsize=(5.5, 3))
    fig.patch.set_facecolor('white')
    ax.set_facecolor('white')
    ax.axis('off')

    # Flow boxes
    boxes = [
        (0.02, 0.6, 0.16, 0.3, "Input\nImage", "#0D6EB8"),
        (0.22, 0.6, 0.16, 0.3, "Sakaguchi\nPreprocessing", "#1565C0"),
        (0.42, 0.6, 0.16, 0.3, "Multi-Input\nFeatures", "#1976D2"),
        (0.62, 0.6, 0.16, 0.3, "4-Model\nEnsemble", "#2196F3"),
        (0.82, 0.6, 0.16, 0.3, "Disease\nPrediction", "#4CAF50"),
    ]
    for (x, y, w, h, label, color) in boxes:
        rect = plt.Rectangle((x, y), w, h, facecolor=color, edgecolor='white',
                              linewidth=1.5, transform=ax.transAxes)
        ax.add_patch(rect)
        ax.text(x + w/2, y + h/2, label, transform=ax.transAxes,
                ha='center', va='center', color='white', fontsize=8, fontweight='bold')
        if x < 0.82:
            ax.annotate('', xy=(x + w + 0.03, y + h/2), xytext=(x + w, y + h/2),
                        xycoords='axes fraction', textcoords='axes fraction',
                        arrowprops=dict(arrowstyle='->', color='#333', lw=1.5))

    # Tensor details below
    tensor_boxes = [
        (0.05, 0.15, 0.18, 0.25, "224×224\nImage", "#42A5F5"),
        (0.28, 0.15, 0.14, 0.25, "8×8\nTensor", "#64B5F6"),
        (0.46, 0.15, 0.14, 0.25, "12×12\nTensor", "#64B5F6"),
        (0.64, 0.15, 0.14, 0.25, "16×16\nTensor", "#64B5F6"),
    ]
    for (x, y, w, h, label, color) in tensor_boxes:
        rect = plt.Rectangle((x, y), w, h, facecolor=color, edgecolor='white',
                              linewidth=1, alpha=0.8, transform=ax.transAxes)
        ax.add_patch(rect)
        ax.text(x + w/2, y + h/2, label, transform=ax.transAxes,
                ha='center', va='center', color='white', fontsize=7, fontweight='bold')

    ax.text(0.5, 0.05, "Multi-Input Features fed to all 4 models simultaneously",
            transform=ax.transAxes, ha='center', va='center', fontsize=8,
            color='#555', style='italic')

    plt.tight_layout()
    path = os.path.join(ASSETS_DIR, "poster_architecture.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


# ============================================================
# Poster Builder
# ============================================================
def build_poster():
    prs = Presentation()
    prs.slide_width = POSTER_W
    prs.slide_height = POSTER_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # --------------------------------------------------------
    # HEADER — Blue banner
    # --------------------------------------------------------
    header_h = Inches(1.6)
    add_shape(slide, Inches(0), Inches(0), POSTER_W, header_h, VIT_BLUE)

    # VIT logo area (white box placeholder)
    add_shape(slide, Inches(0.3), Inches(0.15), Inches(2.5), Inches(1.3), WHITE)
    tf_logo = add_textbox(slide, Inches(0.4), Inches(0.3), Inches(2.3), Inches(1.0),
                          "VIT", font_size=36, bold=True, color=VIT_BLUE, alignment=PP_ALIGN.CENTER)
    add_para(tf_logo, "Vellore Institute of Technology", font_size=9, bold=True,
             color=VIT_DARK, alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, Inches(3), Inches(0.15), Inches(14.5), Inches(0.7),
                "MULTI-INPUT FRUIT & LEAF DISEASE DETECTION SYSTEM",
                font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Student info
    add_textbox(slide, Inches(3), Inches(0.75), Inches(14.5), Inches(0.35),
                f"{STUDENT_NAME} ({REGISTER_NO})",
                font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Guide name
    add_textbox(slide, Inches(3), Inches(1.05), Inches(14.5), Inches(0.25),
                f"Guide: {GUIDE_NAME}",
                font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

    # School name
    add_textbox(slide, Inches(3), Inches(1.25), Inches(14.5), Inches(0.25),
                SCHOOL_NAME,
                font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # --------------------------------------------------------
    # GitHub link row
    # --------------------------------------------------------
    gh_y = Inches(1.6)
    add_shape(slide, Inches(0), gh_y, POSTER_W, Inches(0.4), LIGHT_GRAY)
    add_textbox(slide, Inches(0.5), gh_y, POSTER_W, Inches(0.4),
                f"GITHUB LINK: {GITHUB_LINK}",
                font_size=11, bold=True, color=HEADING_BLUE, alignment=PP_ALIGN.LEFT)

    # --------------------------------------------------------
    # 3-Column Layout
    # --------------------------------------------------------
    content_y = Inches(2.15)
    content_h = Inches(12.6)
    col_w = Inches(6.2)
    gap = Inches(0.2)
    margin = Inches(0.4)

    col1_x = margin
    col2_x = col1_x + col_w + gap
    col3_x = col2_x + col_w + gap

    # Column borders
    add_bordered_box(slide, col1_x, content_y, col_w, content_h)
    add_bordered_box(slide, col2_x, content_y, col_w, content_h)
    add_bordered_box(slide, col3_x, content_y, col_w, content_h)

    inner_margin = Inches(0.25)
    text_w = col_w - 2 * inner_margin

    # ========================================================
    # COLUMN 1: Introduction, Objectives, Scope
    # ========================================================
    y = content_y + Inches(0.2)

    # --- INTRODUCTION ---
    add_section_heading(slide, col1_x + inner_margin, y, text_w, "INTRODUCTION")
    y += Inches(0.45)

    tf = add_textbox(slide, col1_x + inner_margin, y, text_w, Inches(3.0),
                     "Plant diseases are a major threat to agricultural productivity worldwide, "
                     "causing up to 30% crop losses annually. Early and accurate detection of "
                     "diseases in fruits and leaves is critical for preventing spread and reducing "
                     "economic losses.",
                     font_size=11, color=DARK_GRAY)
    add_para(tf, "", font_size=6)
    add_para(tf, "Traditional manual inspection by agricultural experts is slow, subjective, "
             "and not scalable. With the advancement of deep learning and computer vision, "
             "automated disease detection systems can provide real-time, accurate diagnosis "
             "from simple photographs of plant leaves and fruits.", font_size=11, color=DARK_GRAY)
    add_para(tf, "", font_size=6)
    add_para(tf, "This project develops a multi-input, multi-model deep learning system "
             "that leverages a novel Sakaguchi-based preprocessing pipeline to extract "
             "multi-resolution tensor features, enabling robust disease classification "
             "across 30 disease categories.", font_size=11, color=DARK_GRAY)

    y += Inches(3.8)

    # --- OBJECTIVES ---
    add_section_heading(slide, col1_x + inner_margin, y, text_w, "OBJECTIVES")
    y += Inches(0.45)

    tf2 = add_textbox(slide, col1_x + inner_margin, y, text_w, Inches(3.5),
                      "• Design and implement a Sakaguchi-based preprocessing pipeline "
                      "with multi-resolution tensor conversion (8×8, 12×12, 16×16).",
                      font_size=11, color=DARK_GRAY)
    add_para(tf2, "", font_size=4)
    add_para(tf2, "• Build and compare four distinct classification models: "
             "CNN (MobileNetV2), ANN (MLP), ResNet50, and SVM.", font_size=11, color=DARK_GRAY)
    add_para(tf2, "", font_size=4)
    add_para(tf2, "• Achieve 85%+ validation accuracy on a combined dataset of "
             "82,000+ images across 30 disease classes.", font_size=11, color=DARK_GRAY)
    add_para(tf2, "", font_size=4)
    add_para(tf2, "• Develop an interactive Streamlit web application for real-time "
             "disease diagnosis with multi-image upload support.", font_size=11, color=DARK_GRAY)
    add_para(tf2, "", font_size=4)
    add_para(tf2, "• Evaluate model performance using accuracy, loss curves, "
             "and confidence-based predictions.", font_size=11, color=DARK_GRAY)

    y += Inches(3.8)

    # --- SCOPE OF THE PROJECT ---
    add_section_heading(slide, col1_x + inner_margin, y, text_w, "SCOPE OF THE PROJECT")
    y += Inches(0.45)

    tf3 = add_textbox(slide, col1_x + inner_margin, y, text_w, Inches(3.5),
                      "• Applicable to precision agriculture for early disease intervention "
                      "in crops such as Tomato, Pepper, Potato, Apple, Banana, Grape, "
                      "Mango, and Orange.",
                      font_size=11, color=DARK_GRAY)
    add_para(tf3, "", font_size=4)
    add_para(tf3, "• Can be deployed as a mobile-friendly web tool for farmers "
             "and agricultural extension workers.", font_size=11, color=DARK_GRAY)
    add_para(tf3, "", font_size=4)
    add_para(tf3, "• Extensible to new crop types and diseases by retraining "
             "on additional datasets.", font_size=11, color=DARK_GRAY)
    add_para(tf3, "", font_size=4)
    add_para(tf3, "• Multi-model comparison enables selection of the best model "
             "for deployment based on accuracy vs. speed tradeoffs.", font_size=11, color=DARK_GRAY)

    # ========================================================
    # COLUMN 2: Methodology, Architecture
    # ========================================================
    y = content_y + Inches(0.2)

    # --- METHODOLOGY ---
    add_section_heading(slide, col2_x + inner_margin, y, text_w, "METHODOLOGY")
    y += Inches(0.45)

    tf4 = add_textbox(slide, col2_x + inner_margin, y, text_w, Inches(1.0),
                      "Preprocessing Pipeline (Sakaguchi Method)",
                      font_size=13, bold=True, color=DARK_GRAY)
    add_para(tf4, "The input image undergoes: Resize (224×224) → Gaussian Smoothing "
             "(5×5 kernel) → Normalization [0,1] → Sakaguchi Tensor Conversion "
             "(8×8, 12×12, 16×16 multi-resolution features).", font_size=11, color=DARK_GRAY)
    y += Inches(1.4)

    tf5 = add_textbox(slide, col2_x + inner_margin, y, text_w, Inches(1.5),
                      "Data Augmentation",
                      font_size=13, bold=True, color=DARK_GRAY)
    add_para(tf5, "Applied via Keras ImageDataGenerator during training:", font_size=11, color=DARK_GRAY)
    add_para(tf5, "• Rotation: ±20°  • Zoom: 15%  • Shift: 20%", font_size=10, color=DARK_GRAY)
    add_para(tf5, "• Shear: 15%  • Horizontal Flip: Enabled", font_size=10, color=DARK_GRAY)
    add_para(tf5, "• Train/Validation Split: 80/20", font_size=10, color=DARK_GRAY)

    y += Inches(1.7)

    # --- ARCHITECTURE ---
    add_section_heading(slide, col2_x + inner_margin, y, text_w, "ARCHITECTURE")
    y += Inches(0.45)

    # Architecture diagram
    arch_path = generate_architecture_diagram()
    if arch_path and os.path.exists(arch_path):
        slide.shapes.add_picture(arch_path, col2_x + Inches(0.1), y, text_w + Inches(0.2), Inches(2.8))
    y += Inches(3.0)

    # Dataset details
    tf6 = add_textbox(slide, col2_x + inner_margin, y, text_w, Inches(3.5),
                      "Datasets Used:",
                      font_size=13, bold=True, color=DARK_GRAY)
    add_para(tf6, "", font_size=4)
    add_para(tf6, "PlantVillage Dataset:", font_size=11, bold=True, color=HEADING_BLUE)
    add_para(tf6, "  50,000+ images | 15 classes (Tomato, Pepper, Potato)", font_size=10, color=DARK_GRAY)
    add_para(tf6, "", font_size=4)
    add_para(tf6, "Fruits Disease Dataset:", font_size=11, bold=True, color=HEADING_BLUE)
    add_para(tf6, "  5 fruits × 3 conditions = 15 classes", font_size=10, color=DARK_GRAY)
    add_para(tf6, "  (Apple, Banana, Grape, Mango, Orange)", font_size=10, color=DARK_GRAY)
    add_para(tf6, "", font_size=4)
    add_para(tf6, "Combined: 30 classes | 82,000+ images | 224×224 input", font_size=11,
             bold=True, color=DARK_GRAY)
    add_para(tf6, "", font_size=6)
    add_para(tf6, "Four models are trained and compared:", font_size=11, color=DARK_GRAY)
    add_para(tf6, "1. CNN (MobileNetV2) — Transfer Learning", font_size=10, color=DARK_GRAY)
    add_para(tf6, "2. ANN (MLP) — Fully Connected Network", font_size=10, color=DARK_GRAY)
    add_para(tf6, "3. ResNet50 — Deep Residual Network", font_size=10, color=DARK_GRAY)
    add_para(tf6, "4. SVM — Feature Extraction + RBF Classifier", font_size=10, color=DARK_GRAY)

    # ========================================================
    # COLUMN 3: Results, Conclusion, References
    # ========================================================
    y = content_y + Inches(0.2)

    # --- RESULTS AND DISCUSSION ---
    add_section_heading(slide, col3_x + inner_margin, y, text_w, "RESULTS AND DISCUSSION")
    y += Inches(0.45)

    tf7 = add_textbox(slide, col3_x + inner_margin, y, text_w, Inches(0.8),
                      "Model accuracy comparison across all four architectures. "
                      "CNN (MobileNetV2) achieves the highest validation accuracy "
                      "at 85.1%, followed by SVM at 82.3%.",
                      font_size=11, color=DARK_GRAY)
    y += Inches(0.9)

    # Donut charts
    donut_path = generate_model_accuracy_donut()
    if donut_path and os.path.exists(donut_path):
        slide.shapes.add_picture(donut_path, col3_x + Inches(0.5), y, Inches(5.2), Inches(2.5))
    y += Inches(2.7)

    # Caption
    add_textbox(slide, col3_x + inner_margin, y, text_w, Inches(0.3),
                "Fig 1: Validation accuracy of best performing models",
                font_size=9, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    y += Inches(0.5)

    # Bar comparison chart
    bar_path = generate_model_comparison_bars()
    if bar_path and os.path.exists(bar_path):
        slide.shapes.add_picture(bar_path, col3_x + Inches(0.1), y, text_w + Inches(0.2), Inches(2.5))
    y += Inches(2.7)

    # Caption
    add_textbox(slide, col3_x + inner_margin, y, text_w, Inches(0.3),
                "Fig 2: Train vs Validation accuracy comparison across all models",
                font_size=9, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    y += Inches(0.6)

    # --- CONCLUSION ---
    add_section_heading(slide, col3_x + inner_margin, y, text_w, "CONCLUSION")
    y += Inches(0.45)

    tf8 = add_textbox(slide, col3_x + inner_margin, y, text_w, Inches(2.5),
                      "The multi-input architecture with Sakaguchi tensor-based "
                      "preprocessing demonstrates effective disease classification "
                      "across 30 classes of fruits and vegetable leaves.",
                      font_size=11, color=DARK_GRAY)
    add_para(tf8, "", font_size=4)
    add_para(tf8, "CNN (MobileNetV2) with transfer learning achieves the best "
             "validation accuracy of 85.1%, confirming that pre-trained features "
             "combined with multi-resolution inputs yield robust classification.",
             font_size=11, color=DARK_GRAY)
    add_para(tf8, "", font_size=4)
    add_para(tf8, "The Streamlit web application provides a practical, real-time "
             "tool for agricultural disease diagnosis.",
             font_size=11, color=DARK_GRAY)

    y += Inches(2.5)

    # --- REFERENCES ---
    add_section_heading(slide, col3_x + inner_margin, y, text_w, "REFERENCES")
    y += Inches(0.45)

    tf9 = add_textbox(slide, col3_x + inner_margin, y, text_w, Inches(1.5),
                      "[1] PlantVillage Dataset — Kaggle (kaggle.com/datasets/emmarex/plantdisease)",
                      font_size=10, color=DARK_GRAY)
    add_para(tf9, "[2] Fruits Disease Dataset — Kaggle (kaggle.com/datasets/sriramr/fruits-diseases)",
             font_size=10, color=DARK_GRAY)
    add_para(tf9, "[3] MobileNetV2: Inverted Residuals and Linear Bottlenecks — Sandler et al., 2018",
             font_size=10, color=DARK_GRAY)

    # Save
    prs.save(OUTPUT_FILE)
    print(f"\n✅ Poster saved: {OUTPUT_FILE}")


if __name__ == "__main__":
    print("Generating poster...")
    build_poster()
