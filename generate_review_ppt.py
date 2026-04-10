"""
Generate a professional PowerPoint presentation for DS Project Review.
Slides:
  1. Title (Project title, Name, Register No.)
  2. Abstract & Introduction
  3. Dataset Details
  4. EDA Code & Output
  5. Data Pre-Processing & EDA Output
  6. Methodology Planned & Expected Prediction Output
"""

import os
import json
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# CONFIGURATION — UPDATE THESE
# ============================================================
STUDENT_NAME = "Syed Azaan Hussain"   # <-- Your name
REGISTER_NO  = "24BCE5025"
OUTPUT_FILE  = "DS_Project_Review_v2.pptx"

# Color palette
BG_DARK   = RGBColor(0x1B, 0x1B, 0x2F)
BG_CARD   = RGBColor(0x27, 0x27, 0x4A)
ACCENT    = RGBColor(0x00, 0xD2, 0xFF)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xCC, 0xCC)
GREEN     = RGBColor(0x4C, 0xAF, 0x50)
ORANGE    = RGBColor(0xFF, 0x98, 0x00)
RED       = RGBColor(0xFF, 0x4B, 0x4B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ASSETS_DIR = "ppt_assets"
os.makedirs(ASSETS_DIR, exist_ok=True)

# ============================================================
# Helpers
# ============================================================
def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=16, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.font.name = "Calibri"
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ============================================================
# Chart Generators
# ============================================================
def generate_class_distribution_chart():
    """Generate class distribution bar chart from class_names.txt"""
    class_names = []
    if os.path.exists("class_names.txt"):
        with open("class_names.txt", "r") as f:
            class_names = [l.strip() for l in f if l.strip()]

    # Categorize
    categories = {}
    for c in class_names:
        if "Tomato" in c:
            cat = "Tomato"
        elif "Pepper" in c:
            cat = "Pepper"
        elif "Potato" in c:
            cat = "Potato"
        elif "Apple" in c:
            cat = "Apple"
        elif "Banana" in c:
            cat = "Banana"
        elif "Grape" in c:
            cat = "Grape"
        elif "Mango" in c:
            cat = "Mango"
        elif "Orange" in c:
            cat = "Orange"
        else:
            cat = "Other"
        categories[cat] = categories.get(cat, 0) + 1

    fig, ax = plt.subplots(figsize=(8, 4))
    fig.patch.set_facecolor('#1B1B2F')
    ax.set_facecolor('#1B1B2F')
    cats = list(categories.keys())
    vals = list(categories.values())
    colors_list = plt.cm.viridis(np.linspace(0.3, 0.9, len(cats)))
    bars = ax.barh(cats, vals, color=colors_list, edgecolor='none', height=0.6)
    for bar, val in zip(bars, vals):
        ax.text(bar.get_width() + 0.1, bar.get_y() + bar.get_height()/2,
                str(val), va='center', color='white', fontsize=11, fontweight='bold')
    ax.set_xlabel("Number of Classes", color='white', fontsize=12)
    ax.set_title("Classes per Category", color='#00D2FF', fontsize=14, fontweight='bold')
    ax.tick_params(colors='white')
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.xaxis.set_tick_params(labelcolor='white')
    ax.yaxis.set_tick_params(labelcolor='white')
    plt.tight_layout()
    path = os.path.join(ASSETS_DIR, "class_distribution.png")
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor=fig.get_facecolor())
    plt.close(fig)
    return path


def generate_model_comparison_chart():
    """Generate model accuracy comparison from comparison_results.json"""
    results_path = "models/comparison_results.json"
    if not os.path.exists(results_path):
        return None
    with open(results_path, "r") as f:
        data = json.load(f)

    names = [d['model_name'] for d in data]
    train_acc = [d['train_accuracy']*100 for d in data]
    val_acc = [d['val_accuracy']*100 for d in data]

    x = np.arange(len(names))
    width = 0.35

    fig, ax = plt.subplots(figsize=(8, 4))
    fig.patch.set_facecolor('#1B1B2F')
    ax.set_facecolor('#1B1B2F')

    bars1 = ax.bar(x - width/2, train_acc, width, label='Train Acc', color='#00D2FF', edgecolor='none')
    bars2 = ax.bar(x + width/2, val_acc, width, label='Val Acc', color='#4CAF50', edgecolor='none')

    for bar in bars1:
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                f'{bar.get_height():.1f}%', ha='center', va='bottom', color='white', fontsize=9)
    for bar in bars2:
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                f'{bar.get_height():.1f}%', ha='center', va='bottom', color='white', fontsize=9)

    ax.set_ylabel('Accuracy (%)', color='white', fontsize=12)
    ax.set_title('Model Comparison', color='#00D2FF', fontsize=14, fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(names, color='white', fontsize=10)
    ax.legend(facecolor='#27274A', edgecolor='none', labelcolor='white')
    ax.tick_params(colors='white')
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.set_ylim(0, 110)
    plt.tight_layout()
    path = os.path.join(ASSETS_DIR, "model_comparison.png")
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor=fig.get_facecolor())
    plt.close(fig)
    return path


def generate_training_curves_chart():
    """Generate training curves from best model (CNN)"""
    results_path = "models/comparison_results.json"
    if not os.path.exists(results_path):
        return None
    with open(results_path, "r") as f:
        data = json.load(f)

    # Find CNN data
    cnn = next((d for d in data if "CNN" in d['model_name']), data[0])
    history = cnn['history']

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4))
    fig.patch.set_facecolor('#1B1B2F')
    for ax in [ax1, ax2]:
        ax.set_facecolor('#1B1B2F')
        for spine in ax.spines.values():
            spine.set_visible(False)
        ax.tick_params(colors='white')

    epochs = range(1, len(history['accuracy']) + 1)

    ax1.plot(epochs, history['accuracy'], '-o', color='#00D2FF', label='Train', markersize=4)
    ax1.plot(epochs, history['val_accuracy'], '-s', color='#4CAF50', label='Val', markersize=4)
    ax1.set_title('Accuracy', color='#00D2FF', fontsize=13, fontweight='bold')
    ax1.set_xlabel('Epoch', color='white')
    ax1.set_ylabel('Accuracy', color='white')
    ax1.legend(facecolor='#27274A', edgecolor='none', labelcolor='white')

    ax2.plot(epochs, history['loss'], '-o', color='#FF9800', label='Train', markersize=4)
    ax2.plot(epochs, history['val_loss'], '-s', color='#FF4B4B', label='Val', markersize=4)
    ax2.set_title('Loss', color='#FF9800', fontsize=13, fontweight='bold')
    ax2.set_xlabel('Epoch', color='white')
    ax2.set_ylabel('Loss', color='white')
    ax2.legend(facecolor='#27274A', edgecolor='none', labelcolor='white')

    plt.tight_layout()
    path = os.path.join(ASSETS_DIR, "training_curves.png")
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor=fig.get_facecolor())
    plt.close(fig)
    return path


def generate_preprocessing_pipeline_chart():
    """Illustrate the Sakaguchi preprocessing pipeline."""
    fig, ax = plt.subplots(figsize=(10, 3))
    fig.patch.set_facecolor('#1B1B2F')
    ax.set_facecolor('#1B1B2F')
    ax.axis('off')

    steps = [
        ("Load Image\n(RGB)", "#E91E63"),
        ("Resize\n(224x224)", "#9C27B0"),
        ("Gaussian\nSmoothing", "#3F51B5"),
        ("Normalize\n[0, 1]", "#00BCD4"),
        ("Sakaguchi\nTensors\n(8x8,12x12,16x16)", "#4CAF50"),
    ]

    for i, (label, color) in enumerate(steps):
        x = 0.1 + i * 0.18
        rect = plt.Rectangle((x, 0.25), 0.14, 0.5, facecolor=color, edgecolor='white',
                              linewidth=1.5, alpha=0.85, transform=ax.transAxes)
        ax.add_patch(rect)
        ax.text(x + 0.07, 0.5, label, transform=ax.transAxes,
                ha='center', va='center', color='white', fontsize=9, fontweight='bold')

        if i < len(steps) - 1:
            ax.annotate('', xy=(x + 0.17, 0.5), xytext=(x + 0.14, 0.5),
                        xycoords='axes fraction', textcoords='axes fraction',
                        arrowprops=dict(arrowstyle='->', color='white', lw=2))

    ax.set_title("Sakaguchi Preprocessing Pipeline", color='#00D2FF',
                  fontsize=14, fontweight='bold', pad=15)
    plt.tight_layout()
    path = os.path.join(ASSETS_DIR, "preprocessing_pipeline.png")
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor=fig.get_facecolor())
    plt.close(fig)
    return path


def generate_methodology_diagram():
    """Architecture overview diagram."""
    fig, ax = plt.subplots(figsize=(10, 5))
    fig.patch.set_facecolor('#1B1B2F')
    ax.set_facecolor('#1B1B2F')
    ax.axis('off')

    # Input layer
    ax.add_patch(plt.Rectangle((0.02, 0.65), 0.15, 0.25, facecolor='#E91E63',
                                edgecolor='white', lw=1.5, alpha=0.85, transform=ax.transAxes))
    ax.text(0.095, 0.775, "Input\nImage", transform=ax.transAxes,
            ha='center', va='center', color='white', fontsize=10, fontweight='bold')

    # Preprocessing
    ax.add_patch(plt.Rectangle((0.22, 0.65), 0.15, 0.25, facecolor='#9C27B0',
                                edgecolor='white', lw=1.5, alpha=0.85, transform=ax.transAxes))
    ax.text(0.295, 0.775, "Sakaguchi\nPreprocessing", transform=ax.transAxes,
            ha='center', va='center', color='white', fontsize=9, fontweight='bold')

    # Arrow
    ax.annotate('', xy=(0.22, 0.775), xytext=(0.17, 0.775),
                xycoords='axes fraction', textcoords='axes fraction',
                arrowprops=dict(arrowstyle='->', color='white', lw=2))

    # Multi-input split
    ax.annotate('', xy=(0.42, 0.85), xytext=(0.37, 0.775),
                xycoords='axes fraction', textcoords='axes fraction',
                arrowprops=dict(arrowstyle='->', color='white', lw=1.5))
    ax.annotate('', xy=(0.42, 0.65), xytext=(0.37, 0.775),
                xycoords='axes fraction', textcoords='axes fraction',
                arrowprops=dict(arrowstyle='->', color='white', lw=1.5))

    # Main image path
    ax.add_patch(plt.Rectangle((0.42, 0.75), 0.12, 0.18, facecolor='#3F51B5',
                                edgecolor='white', lw=1.5, alpha=0.85, transform=ax.transAxes))
    ax.text(0.48, 0.84, "224x224\nImage", transform=ax.transAxes,
            ha='center', va='center', color='white', fontsize=8, fontweight='bold')

    # Tensor path
    ax.add_patch(plt.Rectangle((0.42, 0.52), 0.12, 0.18, facecolor='#00BCD4',
                                edgecolor='white', lw=1.5, alpha=0.85, transform=ax.transAxes))
    ax.text(0.48, 0.61, "Sakaguchi\nTensors", transform=ax.transAxes,
            ha='center', va='center', color='white', fontsize=8, fontweight='bold')

    # Models
    models = ["CNN\n(MobileNetV2)", "ANN\n(MLP)", "ResNet50", "SVM"]
    colors = ["#4CAF50", "#FF9800", "#2196F3", "#F44336"]
    for i, (name, color) in enumerate(zip(models, colors)):
        y = 0.1 + i * 0.22
        ax.add_patch(plt.Rectangle((0.62, y), 0.14, 0.17, facecolor=color,
                                    edgecolor='white', lw=1.5, alpha=0.85, transform=ax.transAxes))
        ax.text(0.69, y + 0.085, name, transform=ax.transAxes,
                ha='center', va='center', color='white', fontsize=8, fontweight='bold')

    # Arrows to models
    for i in range(4):
        y = 0.1 + i * 0.22 + 0.085
        ax.annotate('', xy=(0.62, y), xytext=(0.56, 0.775),
                    xycoords='axes fraction', textcoords='axes fraction',
                    arrowprops=dict(arrowstyle='->', color='white', lw=1, alpha=0.6))

    # Output
    ax.add_patch(plt.Rectangle((0.82, 0.35), 0.16, 0.3, facecolor='#FFD700',
                                edgecolor='white', lw=2, alpha=0.9, transform=ax.transAxes))
    ax.text(0.90, 0.5, "System\nConsensus\nPrediction", transform=ax.transAxes,
            ha='center', va='center', color='#1B1B2F', fontsize=10, fontweight='bold')

    for i in range(4):
        y = 0.1 + i * 0.22 + 0.085
        ax.annotate('', xy=(0.82, 0.5), xytext=(0.76, y),
                    xycoords='axes fraction', textcoords='axes fraction',
                    arrowprops=dict(arrowstyle='->', color='white', lw=1, alpha=0.6))

    ax.set_title("Multi-Input Multi-Model Architecture", color='#00D2FF',
                  fontsize=14, fontweight='bold', pad=15)
    plt.tight_layout()
    path = os.path.join(ASSETS_DIR, "methodology_diagram.png")
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor=fig.get_facecolor())
    plt.close(fig)
    return path


# ============================================================
# Slide Builders
# ============================================================
def build_slide_1_title(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    # Accent line at top
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W)

    # Project Title
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "Multi-Input Fruit & Leaf Disease Detection System",
                 font_size=36, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.7),
                 "Automated Deep Learning System using Sakaguchi-based Preprocessing",
                 font_size=20, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Accent line
    add_accent_line(slide, Inches(4), Inches(3.8), Inches(5))

    # Name & Reg No
    add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.6),
                 f"Name: {STUDENT_NAME}",
                 font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
                 f"Register No: {REGISTER_NO}",
                 font_size=20, color=LIGHT, alignment=PP_ALIGN.CENTER)

    # Bottom tag
    add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
                 "DS Project Review Presentation",
                 font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)


def build_slide_2_abstract(prs):
    """Slide 2: Abstract & Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Abstract & Introduction", font_size=30, bold=True, color=ACCENT)

    # Abstract card
    add_rounded_rect(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(2.5))
    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.3),
                      "Abstract", font_size=20, bold=True, color=ACCENT)
    add_paragraph(tf,
        "This project implements an automated fruit and leaf disease detection system using "
        "a multi-input, multi-model deep learning architecture. The system processes images "
        "through a Sakaguchi-based preprocessing pipeline that generates multi-resolution "
        "tensors (8x8, 12x12, 16x16) alongside the primary 224x224 image input. Four "
        "distinct models — CNN (MobileNetV2), ANN (MLP), ResNet50, and SVM — are trained "
        "and compared. The best-performing model achieves 85%+ validation accuracy across "
        "30 disease classes spanning fruits and vegetable leaves.",
        font_size=14, color=LIGHT)

    # Introduction card
    add_rounded_rect(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(3.0))
    tf2 = add_text_box(slide, Inches(0.8), Inches(4.1), Inches(11.5), Inches(2.8),
                       "Introduction", font_size=20, bold=True, color=ACCENT)
    add_paragraph(tf2,
        "• Plant diseases cause up to 30% crop losses globally, threatening food security.",
        font_size=14, color=LIGHT)
    add_paragraph(tf2,
        "• Early and accurate detection can prevent spread and reduce economic losses.",
        font_size=14, color=LIGHT)
    add_paragraph(tf2,
        "• Traditional manual inspection is slow, subjective, and requires expert knowledge.",
        font_size=14, color=LIGHT)
    add_paragraph(tf2,
        "• Deep learning-based image classification offers automated, scalable, real-time detection.",
        font_size=14, color=LIGHT)
    add_paragraph(tf2,
        "• Our system combines multiple deep learning architectures with a novel Sakaguchi "
        "tensor-based preprocessing pipeline for robust multi-input classification.",
        font_size=14, color=LIGHT)


def build_slide_3_dataset(prs):
    """Slide 3: Dataset Details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Dataset Details", font_size=30, bold=True, color=ACCENT)

    # Dataset 1
    add_rounded_rect(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.8))
    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.3), Inches(2.6),
                      "PlantVillage Dataset", font_size=20, bold=True, color=GREEN)
    add_paragraph(tf, "• 50,000+ images of healthy & diseased plant leaves", font_size=13, color=LIGHT)
    add_paragraph(tf, "• 15 classes: Tomato, Pepper, Potato diseases", font_size=13, color=LIGHT)
    add_paragraph(tf, "• Source: Kaggle (PlantVillage)", font_size=13, color=LIGHT)
    add_paragraph(tf, "• Format: JPEG/PNG, various resolutions", font_size=13, color=LIGHT)
    add_paragraph(tf, "• Used for training vegetable leaf disease detection", font_size=13, color=LIGHT)

    # Dataset 2
    add_rounded_rect(slide, Inches(6.7), Inches(1.2), Inches(5.8), Inches(2.8))
    tf2 = add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.3), Inches(2.6),
                       "Fruits Disease Dataset", font_size=20, bold=True, color=ORANGE)
    add_paragraph(tf2, "• 5 fruit categories: Apple, Banana, Grape, Mango, Orange", font_size=13, color=LIGHT)
    add_paragraph(tf2, "• 15 classes: Fresh, Rotten, Formalin-mixed per fruit", font_size=13, color=LIGHT)
    add_paragraph(tf2, "• Source: Kaggle (Fruits Disease Dataset)", font_size=13, color=LIGHT)
    add_paragraph(tf2, "• Labeled as Healthy vs. Diseased conditions", font_size=13, color=LIGHT)
    add_paragraph(tf2, "• Used for training fruit disease detection", font_size=13, color=LIGHT)

    # Combined stats card
    add_rounded_rect(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(1.2))
    tf3 = add_text_box(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(1.0),
                       "Combined Dataset Statistics", font_size=18, bold=True, color=ACCENT)
    add_paragraph(tf3,
        "Total Classes: 30  |  Total Images: ~82,000+  |  Train/Val Split: 80/20  |  Image Size: 224x224",
        font_size=14, color=LIGHT)

    # Class distribution chart
    chart_path = generate_class_distribution_chart()
    if chart_path and os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, Inches(2.5), Inches(5.6), Inches(8), Inches(1.7))


def build_slide_4_eda(prs):
    """Slide 4: EDA Code & Output"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Exploratory Data Analysis (EDA)", font_size=30, bold=True, color=ACCENT)

    # EDA Code card
    add_rounded_rect(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5))
    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
                      "EDA Code Snippet", font_size=18, bold=True, color=ACCENT)

    code_text = (
        "# Load and explore the dataset\n"
        "import os\n"
        "import matplotlib.pyplot as plt\n\n"
        "data_dir = 'data/training_data'\n"
        "classes = sorted(os.listdir(data_dir))\n"
        "print(f'Total classes: {len(classes)}')\n\n"
        "# Count images per class\n"
        "counts = {}\n"
        "for cls in classes:\n"
        "    path = os.path.join(data_dir, cls)\n"
        "    counts[cls] = len(os.listdir(path))\n\n"
        "# Visualize class distribution\n"
        "plt.barh(list(counts.keys()),\n"
        "         list(counts.values()))\n"
        "plt.title('Class Distribution')\n"
        "plt.xlabel('Number of Images')\n"
        "plt.show()"
    )
    add_paragraph(tf, code_text, font_size=10, color=LIGHT)

    # EDA Output card
    add_rounded_rect(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5))
    tf2 = add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.5), Inches(0.5),
                       "EDA Output & Findings", font_size=18, bold=True, color=ACCENT)
    add_paragraph(tf2, "Key Findings:", font_size=15, bold=True, color=WHITE)
    add_paragraph(tf2, "• Total Classes: 30 (15 PlantVillage + 15 Fruits)", font_size=12, color=LIGHT)
    add_paragraph(tf2, "• Total Images: ~82,000+", font_size=12, color=LIGHT)
    add_paragraph(tf2, "• Tomato has the most classes (10 diseases)", font_size=12, color=LIGHT)
    add_paragraph(tf2, "• Each fruit has 3 conditions (Fresh/Rotten/Formalin)", font_size=12, color=LIGHT)
    add_paragraph(tf2, "• Dataset is moderately balanced across categories", font_size=12, color=LIGHT)
    add_paragraph(tf2, "", font_size=8, color=LIGHT)
    add_paragraph(tf2, "Class Distribution by Category:", font_size=14, bold=True, color=WHITE)
    add_paragraph(tf2, "  Tomato:   10 classes  |  Pepper:  2 classes", font_size=12, color=LIGHT)
    add_paragraph(tf2, "  Potato:    3 classes  |  Apple:   3 classes", font_size=12, color=LIGHT)
    add_paragraph(tf2, "  Banana:    3 classes  |  Grape:   3 classes", font_size=12, color=LIGHT)
    add_paragraph(tf2, "  Mango:     3 classes  |  Orange:  3 classes", font_size=12, color=LIGHT)


def build_slide_5_preprocessing(prs):
    """Slide 5: Data Pre-Processing Pipeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Data Pre-Processing Pipeline", font_size=30, bold=True, color=ACCENT)

    # Pipeline diagram
    pipeline_path = generate_preprocessing_pipeline_chart()
    if pipeline_path and os.path.exists(pipeline_path):
        slide.shapes.add_picture(pipeline_path, Inches(0.5), Inches(1.1), Inches(12), Inches(2.8))

    # Pre-processing code — full slide width
    add_rounded_rect(slide, Inches(0.5), Inches(4.1), Inches(6), Inches(3.1))
    tf = add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.5),
                      "Sakaguchi Preprocessing Code", font_size=18, bold=True, color=ACCENT)

    code = (
        "def preprocess_pipeline(image_file):\n"
        "  # 1. Load Original Image (RGB)\n"
        "  raw = load_image(image_file)\n\n"
        "  # 2. Resize to 224x224\n"
        "  resized = cv2.resize(raw, (224, 224))\n\n"
        "  # 3. Gaussian Blur (Noise Smoothing)\n"
        "  smoothed = cv2.GaussianBlur(resized, (5,5), 0)\n\n"
        "  # 4. Normalize to [0, 1]\n"
        "  normalized = smoothed / 255.0\n\n"
        "  # 5. Sakaguchi Tensor Conversion\n"
        "  tensors = {\n"
        "    '8x8':  cv2.resize(normalized, (8,8)),\n"
        "    '12x12': cv2.resize(normalized, (12,12)),\n"
        "    '16x16': cv2.resize(normalized, (16,16))\n"
        "  }\n"
        "  return normalized, tensors"
    )
    add_paragraph(tf, code, font_size=10, color=LIGHT)

    # Pipeline description box
    add_rounded_rect(slide, Inches(6.8), Inches(4.1), Inches(6), Inches(3.1))
    tf2 = add_text_box(slide, Inches(7.1), Inches(4.2), Inches(5.5), Inches(0.5),
                       "Pipeline Steps Explained", font_size=18, bold=True, color=ACCENT)
    add_paragraph(tf2, "Step 1 — Load Image:", font_size=13, bold=True, color=WHITE)
    add_paragraph(tf2, "  Read image from upload, convert BGR → RGB", font_size=12, color=LIGHT)
    add_paragraph(tf2, "Step 2 — Resize:", font_size=13, bold=True, color=WHITE)
    add_paragraph(tf2, "  Standardize to 224×224 for model input", font_size=12, color=LIGHT)
    add_paragraph(tf2, "Step 3 — Gaussian Smoothing:", font_size=13, bold=True, color=WHITE)
    add_paragraph(tf2, "  5×5 kernel removes noise artifacts", font_size=12, color=LIGHT)
    add_paragraph(tf2, "Step 4 — Normalize:", font_size=13, bold=True, color=WHITE)
    add_paragraph(tf2, "  Scale pixel values from [0,255] → [0,1]", font_size=12, color=LIGHT)
    add_paragraph(tf2, "Step 5 — Sakaguchi Tensors:", font_size=13, bold=True, color=WHITE)
    add_paragraph(tf2, "  Generate 8×8, 12×12, 16×16 multi-scale features", font_size=12, color=LIGHT)


def build_slide_6_augmentation(prs):
    """Slide 6: Data Augmentation & Pre-Processing EDA Output"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Data Augmentation & Pre-Processing EDA Output", font_size=28, bold=True, color=ACCENT)

    # Augmentation card
    add_rounded_rect(slide, Inches(0.5), Inches(1.2), Inches(6), Inches(5.5))
    tf = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
                      "Training-Time Data Augmentation", font_size=20, bold=True, color=ACCENT)
    add_paragraph(tf, "Applied via Keras ImageDataGenerator:", font_size=14, bold=True, color=WHITE)
    add_paragraph(tf, "", font_size=6, color=LIGHT)
    add_paragraph(tf, "• Rotation Range: ±20°", font_size=14, color=LIGHT)
    add_paragraph(tf, "  → Handles images taken at different angles", font_size=11, color=LIGHT)
    add_paragraph(tf, "• Zoom Range: 15%", font_size=14, color=LIGHT)
    add_paragraph(tf, "  → Simulates varying camera distances", font_size=11, color=LIGHT)
    add_paragraph(tf, "• Width/Height Shift: 20%", font_size=14, color=LIGHT)
    add_paragraph(tf, "  → Handles off-center subjects", font_size=11, color=LIGHT)
    add_paragraph(tf, "• Shear Range: 15%", font_size=14, color=LIGHT)
    add_paragraph(tf, "  → Adds perspective distortion", font_size=11, color=LIGHT)
    add_paragraph(tf, "• Horizontal Flip: Enabled", font_size=14, color=LIGHT)
    add_paragraph(tf, "  → Doubles effective dataset size", font_size=11, color=LIGHT)
    add_paragraph(tf, "", font_size=6, color=LIGHT)
    add_paragraph(tf, "Augmentation prevents overfitting and improves", font_size=13, color=GREEN)
    add_paragraph(tf, "generalization on unseen disease images.", font_size=13, color=GREEN)

    # Pre-processing EDA output card
    add_rounded_rect(slide, Inches(6.8), Inches(1.2), Inches(6), Inches(5.5))
    tf2 = add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.5), Inches(0.5),
                       "Pre-Processing EDA Findings", font_size=20, bold=True, color=ACCENT)
    add_paragraph(tf2, "Image Statistics (after preprocessing):", font_size=14, bold=True, color=WHITE)
    add_paragraph(tf2, "", font_size=6, color=LIGHT)
    add_paragraph(tf2, "• Input Resolution: 224 × 224 × 3 (RGB)", font_size=13, color=LIGHT)
    add_paragraph(tf2, "• Pixel Range: [0.0, 1.0] after normalization", font_size=13, color=LIGHT)
    add_paragraph(tf2, "• Gaussian Blur Kernel: 5 × 5", font_size=13, color=LIGHT)
    add_paragraph(tf2, "• Batch Size: 32 images per iteration", font_size=13, color=LIGHT)
    add_paragraph(tf2, "", font_size=6, color=LIGHT)
    add_paragraph(tf2, "Sakaguchi Tensor Output Shapes:", font_size=14, bold=True, color=WHITE)
    add_paragraph(tf2, "  • Tensor 1:  8 × 8 × 3   =    192 features", font_size=13, color=LIGHT)
    add_paragraph(tf2, "  • Tensor 2: 12 × 12 × 3  =    432 features", font_size=13, color=LIGHT)
    add_paragraph(tf2, "  • Tensor 3: 16 × 16 × 3  =    768 features", font_size=13, color=LIGHT)
    add_paragraph(tf2, "", font_size=6, color=LIGHT)
    add_paragraph(tf2, "Key Insight:", font_size=14, bold=True, color=WHITE)
    add_paragraph(tf2, "Resize → Blur order ensures the 5×5 kernel", font_size=12, color=LIGHT)
    add_paragraph(tf2, "has consistent relative effect across all images,", font_size=12, color=LIGHT)
    add_paragraph(tf2, "matching training and inference pipelines.", font_size=12, color=LIGHT)


def build_slide_7_methodology(prs):
    """Slide 7: Methodology & Model Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Methodology & Model Architecture", font_size=30, bold=True, color=ACCENT)

    # Architecture diagram
    methodology_path = generate_methodology_diagram()
    if methodology_path and os.path.exists(methodology_path):
        slide.shapes.add_picture(methodology_path, Inches(0.3), Inches(1.1), Inches(8), Inches(4.0))

    # Model details card
    add_rounded_rect(slide, Inches(8.5), Inches(1.1), Inches(4.5), Inches(6.0))
    tf = add_text_box(slide, Inches(8.8), Inches(1.2), Inches(4.0), Inches(0.5),
                      "4 Models Compared", font_size=20, bold=True, color=ACCENT)

    add_paragraph(tf, "1. CNN (MobileNetV2)", font_size=14, bold=True, color=GREEN)
    add_paragraph(tf, "   Transfer learning, frozen base", font_size=11, color=LIGHT)
    add_paragraph(tf, "   + GlobalAvgPool + Dense(256)", font_size=11, color=LIGHT)
    add_paragraph(tf, "", font_size=4, color=LIGHT)

    add_paragraph(tf, "2. ANN (MLP)", font_size=14, bold=True, color=ORANGE)
    add_paragraph(tf, "   Flatten + Dense(1024→512→256→128)", font_size=11, color=LIGHT)
    add_paragraph(tf, "   BatchNorm + Dropout(0.5)", font_size=11, color=LIGHT)
    add_paragraph(tf, "", font_size=4, color=LIGHT)

    add_paragraph(tf, "3. ResNet50", font_size=14, bold=True, color=ACCENT)
    add_paragraph(tf, "   Fine-tuned top 50 layers", font_size=11, color=LIGHT)
    add_paragraph(tf, "   + Dense(512→256) + BatchNorm", font_size=11, color=LIGHT)
    add_paragraph(tf, "", font_size=4, color=LIGHT)

    add_paragraph(tf, "4. SVM", font_size=14, bold=True, color=RED)
    add_paragraph(tf, "   MobileNetV2 feature extractor", font_size=11, color=LIGHT)
    add_paragraph(tf, "   + RBF kernel SVM (C=10)", font_size=11, color=LIGHT)
    add_paragraph(tf, "", font_size=6, color=LIGHT)

    add_paragraph(tf, "All models use multi-input:", font_size=12, bold=True, color=WHITE)
    add_paragraph(tf, "224×224 image + 3 Sakaguchi tensors", font_size=11, color=GREEN)

    # Training config card
    add_rounded_rect(slide, Inches(0.3), Inches(5.3), Inches(8), Inches(1.8))
    tf2 = add_text_box(slide, Inches(0.6), Inches(5.4), Inches(7.5), Inches(0.5),
                       "Training Configuration", font_size=18, bold=True, color=ACCENT)
    add_paragraph(tf2,
        "Optimizer: Adam  |  Loss: Categorical Cross-Entropy  |  Callbacks: EarlyStopping, ReduceLR, ModelCheckpoint",
        font_size=13, color=LIGHT)
    add_paragraph(tf2,
        "Epochs: 20 (with early stopping)  |  Batch Size: 32  |  Validation Split: 20%",
        font_size=13, color=LIGHT)


def build_slide_8_results(prs):
    """Slide 8: Expected Prediction Output & Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_line(slide, Inches(0), Inches(0), SLIDE_W)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                 "Expected Prediction Output & Conclusion", font_size=28, bold=True, color=ACCENT)

    # Model comparison chart
    comp_path = generate_model_comparison_chart()
    if comp_path and os.path.exists(comp_path):
        slide.shapes.add_picture(comp_path, Inches(0.3), Inches(1.1), Inches(6.5), Inches(3.2))

    # Training curves
    curves_path = generate_training_curves_chart()
    if curves_path and os.path.exists(curves_path):
        slide.shapes.add_picture(curves_path, Inches(0.3), Inches(4.5), Inches(6.5), Inches(2.8))

    # Expected prediction output card
    add_rounded_rect(slide, Inches(7.0), Inches(1.1), Inches(6), Inches(3.0))
    tf = add_text_box(slide, Inches(7.3), Inches(1.2), Inches(5.5), Inches(0.5),
                      "Expected Prediction Output", font_size=20, bold=True, color=ACCENT)
    add_paragraph(tf, "For each uploaded image, the system outputs:", font_size=14, color=WHITE)
    add_paragraph(tf, "", font_size=4, color=LIGHT)
    add_paragraph(tf, "• Disease Name (e.g., Tomato Late Blight)", font_size=13, color=LIGHT)
    add_paragraph(tf, "• Confidence Score (e.g., 92.3%)", font_size=13, color=LIGHT)
    add_paragraph(tf, "• Health Status: Healthy / Diseased", font_size=13, color=LIGHT)
    add_paragraph(tf, "• Top-3 predictions with probability bars", font_size=13, color=LIGHT)
    add_paragraph(tf, "• Sakaguchi tensor metadata", font_size=13, color=LIGHT)
    add_paragraph(tf, "", font_size=4, color=LIGHT)
    add_paragraph(tf, "Best Model: CNN (MobileNetV2) @ 85.1%", font_size=14, bold=True, color=GREEN)

    # Conclusion card
    add_rounded_rect(slide, Inches(7.0), Inches(4.4), Inches(6), Inches(2.9))
    tf2 = add_text_box(slide, Inches(7.3), Inches(4.5), Inches(5.5), Inches(0.5),
                       "Conclusion", font_size=20, bold=True, color=ACCENT)
    add_paragraph(tf2, "• Multi-input architecture with Sakaguchi tensors", font_size=13, color=LIGHT)
    add_paragraph(tf2, "  enhances feature extraction at multiple scales.", font_size=12, color=LIGHT)
    add_paragraph(tf2, "• CNN (MobileNetV2) achieves best accuracy (85%+)", font_size=13, color=LIGHT)
    add_paragraph(tf2, "  with efficient transfer learning.", font_size=12, color=LIGHT)
    add_paragraph(tf2, "• SVM provides a strong non-DL baseline (82%).", font_size=13, color=LIGHT)
    add_paragraph(tf2, "• Streamlit web app enables real-time diagnosis.", font_size=13, color=LIGHT)
    add_paragraph(tf2, "", font_size=6, color=LIGHT)
    add_paragraph(tf2, "Thank You!", font_size=22, bold=True, color=ACCENT,
                  alignment=PP_ALIGN.CENTER)


# ============================================================
# Main
# ============================================================
def main():
    print("Generating presentation assets...")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("  Building Slide 1: Title...")
    build_slide_1_title(prs)

    print("  Building Slide 2: Abstract & Introduction...")
    build_slide_2_abstract(prs)

    print("  Building Slide 3: Dataset Details...")
    build_slide_3_dataset(prs)

    print("  Building Slide 4: EDA Code & Output...")
    build_slide_4_eda(prs)

    print("  Building Slide 5: Data Pre-Processing Pipeline...")
    build_slide_5_preprocessing(prs)

    print("  Building Slide 6: Data Augmentation & Pre-Processing EDA...")
    build_slide_6_augmentation(prs)

    print("  Building Slide 7: Methodology & Model Architecture...")
    build_slide_7_methodology(prs)

    print("  Building Slide 8: Expected Output & Conclusion...")
    build_slide_8_results(prs)

    prs.save(OUTPUT_FILE)
    print(f"\n✅ Presentation saved: {OUTPUT_FILE}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
